import logging
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import google.generativeai as genai
from PIL import Image

def _extract_content_from_shape(shape):
    """
    Extracts both text and image data from a shape and its children.
    Returns a tuple of (list_of_text_strings, list_of_image_objects).
    """
    text_parts = []
    image_parts = []

    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        for row in shape.table.rows:
            row_text = " | ".join([cell.text for cell in row.cells])
            text_parts.append(row_text)
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            sub_text, sub_images = _extract_content_from_shape(sub_shape)
            text_parts.extend(sub_text)
            image_parts.extend(sub_images)
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            image = Image.open(io.BytesIO(shape.image.blob))
            image_parts.append(image)
        except Exception as e:
            logging.warning(f"Could not process an image shape in PPTX: {e}")
    elif hasattr(shape, "text"):
        if shape.text:
            text_parts.append(shape.text)
            
    return text_parts, image_parts

def read_ppt_and_chunk(ppt_path: str):
    """
    Reads a PowerPoint file using a hybrid, multimodal approach with a strict data-extraction prompt.
    """
    try:
        prs = Presentation(ppt_path)
        model = genai.GenerativeModel('gemini-1.5-pro-latest') 

        for i, slide in enumerate(prs.slides):
            slide_text_parts = []
            slide_image_parts = []
            
            for shape in slide.shapes:
                text, images = _extract_content_from_shape(shape)
                slide_text_parts.extend(text)
                slide_image_parts.extend(images)
            
            if not slide_text_parts and not slide_image_parts:
                continue

            # --- THIS IS THE FIX ---
            # Using a strict, rule-based prompt identical in spirit to the image parser.
            prompt_parts = [
                """You are a precision OCR and data extraction engine. Your only task is to transcribe and synthesize the content from the provided slide materials (text and images).

**Instructions:**
1. Integrate the provided machine-readable text and perform OCR on the provided images.
2. Combine all information into a single, coherent block of text representing the slide's content.
3. If you see a table, transcribe its content, using pipes (|) to separate columns.
4. **Crucially, DO NOT add any commentary, explanations, or introductory sentences like 'This slide shows...'.** Do not describe the slide's layout. Only output the transcribed and synthesized data."""
            ]
            
            if slide_text_parts:
                combined_text = "\n".join(slide_text_parts)
                prompt_parts.append("\n--- Machine-Readable Text ---\n" + combined_text)

            prompt_parts.extend(slide_image_parts)

            logging.info(f"Processing Slide {i+1} with multimodal vision...")
            response = model.generate_content(prompt_parts)
            
            if response.text:
                yield {
                    "text": response.text.strip(),
                    "source": f"Slide {i + 1}"
                }

    except Exception as e:
        logging.error(f"Could not process PowerPoint file '{ppt_path}': {e}", exc_info=True)
        return